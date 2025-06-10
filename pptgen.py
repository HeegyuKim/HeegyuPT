import litellm
import json
from pydantic import BaseModel
import re
from tqdm.auto import tqdm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
import os
import asyncio


class SlideRequest(BaseModel):
    title: str
    content: str
    note: str = None  # Optional field for speaker notes

    def to_dict(self):
        return {
            "title": self.title,
            "content": self.content,
            "note": self.note,
        }
    def __str__(self):
        return f"Slide(title={self.title}, content={self.content}, note={self.note})"
    
class PresentationRequest(BaseModel):
    title: str
    slides: list[SlideRequest]

    def to_dict(self):
        return {
            "title": self.title,
            "slides": [slide.to_dict() for slide in self.slides],
        }

    def __str__(self):
        return f"Presentation(title={self.title}, slides={self.slides})"
    
    def to_json(self):
        return json.dumps(self.to_dict(), indent=2)
    
class SectionRequest(BaseModel):
    title: str
    description: str
    num_slides: int
    
class StructureRequest(BaseModel):
    title: str
    description: str
    num_total_slides: int
    sections: list[SectionRequest]


SYSTEM_PROMPT = """
You are a helpful assistant that generates PowerPoint presentations based on user prompts. Each presentation consists of a title and multiple slides, each with a title and content. The content can include text, bullet points, or other relevant information.
Here’s a concise English guideline to produce well-structured, audience-friendly PowerPoint decks:

## 1. Slide Structure & Word-Count  
-  One main idea per slide.  
-  Title: ≤ 8 words.  
-  Body text: ≤ 100-150 words per slide.
-  Bullet list: 3–5 bullets, including descriptive sentences. (e.g., "Findings: we found that...")

## 2. Visual & Layout Best Practices  
-  Leave ample white-space: don’t fill every inch—let visuals “breathe.”  
-  Incorporate one strong visual (photo, chart or icon) per slide.  
-  Apply a consistent color palette (2–3 theme colors) and master slide.  
-  Limit animations to simple builds (fade, appear) and avoid over-animation.  

## 3. Timing & Slide Count  
-  Rule of thumb: 1 slide per 1–2 minutes of talk time.  
-  5 min talk → 4–6 slides  
-  10 min talk → 8–12 slides  
-  20 min talk → 12–18 slides  
-  Schedule 1–2 min at start for title/agenda, and 1–2 min at end for Q&A/summary.  

## 4. Speaker Notes (“Presenter Captions”)  
-  Write in full sentences—this is your script/guide, not the audience’s.  
-  For each slide:  
  – **Intro sentence**: transition from previous slide.  
  – **Key points**: elaborate each bullet (1–2 sentences per bullet).  
  – **Closing/CTA**: summary line or “next” cue.  
"""

async def generate_presentation(prompt: str, model: str = "openai/gpt-4o-search-preview") -> PresentationRequest:
    response = await litellm.acompletion(
        model=model,
        messages=[
            {
                "role": "system",
                "content": SYSTEM_PROMPT,
            },
            {
                "role": "user",
                "content": "PPT Contents:" + prompt,
            }
        ],
        max_tokens=16384,
    )
    print("Response from model:\n", response.choices[0].message.content)
    
    response = await litellm.acompletion(
        model="gpt-4.1-mini",
        messages=[
            {
                "role": "user",
                "content": "Convert the following ppt contents into a structured JSON format for a PowerPoint presentation.:\n\n" + response.choices[0].message.content,
            }
        ],
        max_tokens=16384,
        response_format=PresentationRequest,
    )
    presentation = json.loads(response.choices[0].message.content)
    return PresentationRequest(**presentation)


def add_markdown_to_placeholder(placeholder, md_text):
    # text_frame 초기화
    tf = placeholder.text_frame
    tf.clear()

    # 링크용 정규식: [text](url)
    link_pattern = re.compile(r'\[([^\]]+)\]\(([^)]+)\)')

    for line in md_text.splitlines():
        # 1) Heading: #, ##, ###
        m = re.match(r'^(#{1,6})\s+(.*)', line)
        if m:
            level = len(m.group(1))
            content = m.group(2)
            p = tf.add_paragraph()
            p.text = content
            p.font.bold = True
            # Heading 레벨에 따라 폰트 크기 조정
            p.font.size = Pt(32 - (level-1)*4)
            continue

        # 2) Bullet list: -, *, + 로 시작
        m = re.match(r'^(\s*)([-\*\+])\s+(.*)', line)
        if m:
            indent_spaces = len(m.group(1))
            content = m.group(3)
            # Bullet list 처리 부분 중에서
            p = tf.add_paragraph()
            p.level = indent_spaces // 2   # 들여쓰기 레벨 설정
            # bullet 활성화
            p._p.get_or_add_pPr().bullet = True
            # 폰트 크기 등 스타일 지정
            line = content

        else:
            # 그냥 일반 단락
            p = tf.add_paragraph()
            p.text = ""

        # 1) 링크가 없는 경우 그냥 run으로 삽입
        if not link_pattern.search(line):
            r = p.add_run()
            r.text = line
            continue
        # 2) 링크가 있는 경우 run 단위로 쪼개서 삽입
        else:
            last_end = 0
            for lm in link_pattern.finditer(line):
                # 링크 전 텍스트
                pre_text = line[last_end:lm.start()]
                if pre_text:
                    r = p.add_run()
                    r.text = pre_text

                # 링크 텍스트
                link_text, link_url = lm.group(1), lm.group(2)
                r = p.add_run()
                r.text = link_text
                # 하이퍼링크 속성 지정
                h = r.hyperlink
                h.address = link_url
                # 링크 색상을 PPT 기본 accent로 지정
                r.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                r.font.underline = True

                last_end = lm.end()

            # 4) 링크 뒤 남은 텍스트
            tail = line[last_end:]
            if tail:
                r = p.add_run()
                r.text = tail


def md_to_pptx(presentation: PresentationRequest) -> Presentation:
    # 16:9 비율의 PowerPoint 프레젠테이션 생성
    prs = Presentation()
    # prs.slide_width = Inches(16)
    # prs.slide_height = Inches(9)

    # 제목 슬라이드 추가
    title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드 레이아웃
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = presentation.title

    for slide in presentation.slides:
        layout = prs.slide_layouts[1]  # 제목 + 본문 슬라이드
        ppt_slide = prs.slides.add_slide(layout)
        # 제목
        ppt_slide.shapes.title.text = slide.title

        if slide.note:
            # 스피커 노트 추가
            ppt_slide.notes_slide.notes_text_frame.text = slide.note

        # 본문 placeholder (보통 index=1)
        ph = ppt_slide.placeholders[1]
        # 본문 크기
        # ph.left = Inches(1)
        # ph.top = Inches(1.5)
        # ph.width = Inches(14)
        # ph.height = Inches(6.5)
        add_markdown_to_placeholder(ph, slide.content)
    
    return prs



async def generate_presentation_structure(user_requirements: str, report: str, model: str = "gpt-4.1-mini") -> StructureRequest:
    SYSTEM_PROMPT = """
You are a helpful assistant that generates structured outlines for PowerPoint presentations based on user reports. Each outline consists of a title, description, total number of slides, and sections with titles, descriptions, and slide counts.
"""
    response = await litellm.acompletion(
        model=model,
        messages=[
            {
                "role": "system",
                "content": SYSTEM_PROMPT,
            },
            {
                "role": "user",
                "content": f"Generate a structured presentation outline for the following report:\n\n---\n# Report:\n\n{report}\n\n---\n# User Requirements:\n\n{user_requirements}",
            }
        ],
        max_tokens=16384,
        response_format=StructureRequest,
    )
    return StructureRequest(**json.loads(response.choices[0].message.content))


async def generate_section_slides(user_requirements: str, report: str, section: SectionRequest, model: str = "gpt-4.1-mini") -> list[SlideRequest]:
    SYSTEM_PROMPT = """
You are a helpful assistant that generates PowerPoint slides based on section outlines. Each slide should have a title, content, and optional speaker notes.
Guidelines for generating slides:
- Each slide should focus on a single main idea.
- The title should be concise (≤ 8 words).
- The content should be clear and concise, ideally 100-150 words per slide.
- Use bullet points for clarity, with 3-5 bullets per slide.
"""
    response = await litellm.acompletion(
        model=model,
        messages=[
            {
                "role": "system",
                "content": SYSTEM_PROMPT,
            },
            {
                "role": "user",
                "content": f"Generate {section.num_slides} slides for the section titled '{section.title}' with description '{section.description}' based on the following report.\n\n---\n# Report:\n\n{report}\n\n---\n# User Requirements for the Full Presentation:\n\n{user_requirements}",
            }
        ],
        max_tokens=16384,
        response_format=PresentationRequest
    )
    presentation = PresentationRequest(**json.loads(response.choices[0].message.content))
    return presentation.slides

async def create_presentation_from_report(user_requirements: str, report: str, model: str, filename: str = None) -> str:
    structure = await generate_presentation_structure(user_requirements, report, model=model)
    
    coroutines = [generate_section_slides(user_requirements, report, section, model=model) for section in structure.sections]
    section_slides = await asyncio.gather(*coroutines)
    slides = []
    for i, section_slides in enumerate(section_slides):
        section = structure.sections[i]
        slides.append(SlideRequest(title=f"{i + 1}: {section.title}", content=section.description, note="This is a section overview slide."))
        slides.extend(section_slides)


    presentation = PresentationRequest(title=structure.title, slides=slides)
    # print("Generated Presentation:\n", presentation)
    prs = md_to_pptx(presentation)
    if filename is None:
        os.makedirs("generated_presentations", exist_ok=True)
        filename = f"generated_presentations/{presentation.title.replace(' ', '_').replace('/', '_')}.pptx"
    prs.save(filename)
    return filename


if __name__ == "__main__":
    import asyncio

    # asyncio.run(test_main_generate_from_report())

    # prompt = " https://arxiv.org/abs/2505.24832\n\nGenerate 30 slides for a PowerPoint presentation on this paper"
    # presentation = asyncio.run(generate_presentation(prompt, model="openai/gpt-4o-mini-search-preview-2025-03-11"))
    # print(json.dumps(presentation.to_dict(), indent=2, ensure_ascii=False))
    
    # with open("presentation2.json", "w") as f:
    #     f.write(json.dumps(presentation.to_dict(), indent=2, ensure_ascii=False))
    # with open("presentation2.json", "r") as f:
    #     presentation = PresentationRequest(**json.load(f))

    # # Example of creating a PowerPoint file from the generated presentation
    # prs = md_to_pptx(presentation)
    # filename = presentation.title.replace(" ", "_").replace("/", "_") + ".pptx"
    # prs.save(filename)